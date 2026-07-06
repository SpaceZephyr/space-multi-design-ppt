#!/usr/bin/env python3
"""Export slides to PPTX and/or PDF.

Two sources:
  1. HTML slides (default): renders each NN-*.html at 2x via headless Chromium
     (playwright), then assembles full-bleed 16:9 PPTX / PDF.
  2. --from-images: skips rendering, uses existing NN-*.png in the directory.

Usage:
    python3 export_deck.py <slides_dir> --pptx out.pptx --pdf out.pdf
    python3 export_deck.py <deck_dir> --from-images --pptx out.pptx

Deps: pip install playwright python-pptx img2pdf --break-system-packages
      python3 -m playwright install chromium        (HTML mode only)
"""
import argparse
import re
import sys
import tempfile
from pathlib import Path


def render_html_slides(slides_dir: Path, out_dir: Path, scale: int = 2):
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        sys.exit("playwright not installed. Run: pip install playwright --break-system-packages && python3 -m playwright install chromium")

    files = sorted([p for p in slides_dir.glob("*.html") if re.match(r"^\d+", p.name)], key=lambda p: p.name)
    if not files:
        sys.exit(f"No numbered slide HTML files in {slides_dir}")
    pngs = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=scale)
        for f in files:
            page.goto(f.as_uri())
            page.wait_for_load_state("networkidle")
            page.wait_for_timeout(400)  # let webfonts settle
            png = out_dir / (f.stem + ".png")
            page.screenshot(path=str(png))
            pngs.append(png)
            print(f"rendered {f.name}")
        browser.close()
    return pngs


def collect_images(deck_dir: Path):
    pngs = sorted([p for p in deck_dir.glob("*.png") if re.match(r"^\d+", p.name)], key=lambda p: p.name)
    if not pngs:
        sys.exit(f"No numbered PNG files in {deck_dir}")
    return pngs


def build_pptx(pngs, out_path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        sys.exit("python-pptx not installed. Run: pip install python-pptx --break-system-packages")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(out_path))
    print(f"OK: {out_path} ({len(pngs)} slides)")


def build_pdf(pngs, out_path: Path):
    try:
        import img2pdf
        out_path.write_bytes(img2pdf.convert([str(p) for p in pngs]))
    except ImportError:
        try:
            from PIL import Image
            imgs = [Image.open(p).convert("RGB") for p in pngs]
            imgs[0].save(str(out_path), save_all=True, append_images=imgs[1:])
        except ImportError:
            sys.exit("Neither img2pdf nor Pillow installed. Run: pip install img2pdf --break-system-packages")
    print(f"OK: {out_path} ({len(pngs)} pages)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("source_dir", help="slides/ dir (HTML mode) or deck dir (--from-images)")
    ap.add_argument("--pptx", help="output .pptx path")
    ap.add_argument("--pdf", help="output .pdf path")
    ap.add_argument("--from-images", action="store_true", help="use existing NN-*.png instead of rendering HTML")
    ap.add_argument("--keep-pngs", help="directory to keep rendered PNGs (default: temp)")
    args = ap.parse_args()
    if not args.pptx and not args.pdf:
        sys.exit("Specify --pptx and/or --pdf")

    src = Path(args.source_dir)
    if args.from_images:
        pngs = collect_images(src)
    else:
        out_dir = Path(args.keep_pngs) if args.keep_pngs else Path(tempfile.mkdtemp(prefix="deck-png-"))
        out_dir.mkdir(parents=True, exist_ok=True)
        pngs = render_html_slides(src, out_dir)

    if args.pptx:
        build_pptx(pngs, Path(args.pptx))
    if args.pdf:
        build_pdf(pngs, Path(args.pdf))


if __name__ == "__main__":
    main()
